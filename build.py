#!/usr/bin/env python3
"""
Genera el template de presentacion Perpetual Technologies (.pptx, 16:9).
Reinterpreta el layout "Marketing for your Brand" con el design system Perpetual:
colores de marca, tipografia Armin Grotesk, logos oficiales y reglas duras
(nada de negro puro, fondos oscuros en #0b1220, logo correcto por contraste).

Fuente de verdad de tokens: perpetual-design-system (SKILL.md / tokens.md).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
LOGODIR = os.path.join(HERE, "assets", "logo")
LOGO_COLOR = os.path.join(LOGODIR, "perpetual-color.png")
LOGO_DARK = os.path.join(LOGODIR, "perpetual-dark.png")
OUT = os.path.join(HERE, "dist", "perpetual-marketing-template.pptx")

# ---------------------------------------------------------------------------
# Tokens (de tokens.md)
# ---------------------------------------------------------------------------
def C(hexstr):
    return RGBColor.from_string(hexstr)

BRAND_BLUE = "0032CB"
BRAND_ORANGE = "F33D1F"
BRAND_YELLOW = "FBB900"

BG = "FFFFFF"
SURFACE = "F8F9FC"
SURFACE2 = "EEF1F8"
BORDER = "DDE1EF"
ACCENT = "1A56DB"
ACCENT2 = "F97316"
TEXT = "111827"
TEXT_DIM = "374151"
TEXT_MUTED = "6B7280"

BG_DARK = "0B1220"
SURFACE_DARK = "131C2E"
BORDER_DARK = "243047"
TEXT_ON_DARK = "FFFFFF"
TEXT_DIM_DARK = "9AA6BD"

# Fondo oscuro de SECCIONES del deck (override de marca pedido para este deck:
# el "oscuro" es el azul Perpetual #1a56db, no el casi-negro). Solo deck, no token.
SECTION_DARK = "1A56DB"
SECTION_DARK_DEEP = "123F9E"   # tinte mas oscuro para profundidad sobre el azul
SECTION_DARK_DIM = "C7D7F5"    # texto secundario legible sobre azul

OK = "059669"
WARN = "D97706"
ERROR = "DC2626"
VIOLET = "7E22CE"

# Paleta de datos para charts (orden recomendado en tokens.md)
DATA_PALETTE = [ACCENT, ACCENT2, OK, BRAND_YELLOW, VIOLET, TEXT_MUTED]

# Tipografia Armin Grotesk (instalada en el sistema; familias confirmadas por fc-list)
F_DISPLAY = "Armin Grotesk Black"     # H1 / titulos / numeros grandes (peso maximo con nombre de familia estable)
F_HEAD = "Armin Grotesk SemiBold"     # H2 / subtitulos de seccion
F_LABEL = "Armin Grotesk SemiBold"    # labels / chips uppercase
F_BODY = "Armin Grotesk"              # cuerpo
F_LIGHT = "Armin Grotesk Normal"      # subtitulos ligeros

YEAR = "2026"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C(bg)
    return s


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _soft_shadow(shape, blur=12, dist=3, alpha=88):
    """Sombra suave estilo --shadow-card."""
    sp = shape._element.spPr
    # remove existing effect list
    for tag in ("a:effectLst",):
        el = sp.find(qn(tag))
        if el is not None:
            sp.remove(el)
    eff = sp.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(Emu(Pt(blur)).emu if False else int(blur * 12700)),
        "dist": str(int(dist * 12700)),
        "dir": "5400000",
        "rotWithShape": "0",
    })
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "1A2B5C"})
    a = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - alpha) * 1000))})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    sp.append(eff)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=0.08,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line)
        sp.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        _soft_shadow(sp)
    return sp


def grad_bar(s, x, y, w, h):
    """Barra/elemento con gradiente de marca azul->naranja."""
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_shadow(sp)
    sp.line.fill.background()
    try:
        sp.adjustments[0] = 0.5
    except Exception:
        pass
    # gradiente manual via XML
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})
    for pos, col in ((0, ACCENT), (50000, "3B82F6"), (100000, ACCENT2)):
        gs = gsLst.makeelement(qn("a:gs"), {"pos": str(pos)})
        c = gs.makeelement(qn("a:srgbClr"), {"val": col})
        gs.append(c)
        gsLst.append(gs)
    grad.append(gsLst)
    lin = grad.makeelement(qn("a:lin"), {"ang": "0", "scaled": "1"})
    grad.append(lin)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.insert(list(spPr).index(ln), grad)
    else:
        spPr.append(grad)
    return sp


def text(s, x, y, w, h, runs, size=18, color=TEXT, font=F_BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, line_spacing=None,
         upper=False, wrap=True):
    """runs: str  OR  list de (texto, color, font, size_opt) para runs mixtos."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, color, font, size)]
    for item in runs:
        txt = item[0]
        rcolor = item[1] if len(item) > 1 and item[1] else color
        rfont = item[2] if len(item) > 2 and item[2] else font
        rsize = item[3] if len(item) > 3 and item[3] else size
        r = p.add_run()
        r.text = txt.upper() if upper else txt
        r.font.name = rfont
        r.font.size = Pt(rsize)
        r.font.bold = bold
        r.font.color.rgb = C(rcolor)
        if spacing is not None:
            _letter_spacing(r, spacing)
    return tb


def _letter_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def chip(s, x, y, label, fill=SURFACE2, fg=ACCENT, w=None):
    """Badge/chip uppercase tipo seccion."""
    w = w or (0.16 + 0.085 * len(label))
    sp = rect(s, x, y, w, 0.34, fill=fill, radius=0.5)
    text(s, x, y, w, 0.34, label, size=10.5, color=fg, font=F_LABEL,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2, upper=True)
    return w


def hexagon(s, x, y, size, fill=ACCENT, line=None):
    sp = rect(s, x, y, size, size, fill=fill, line=line, shape=MSO_SHAPE.HEXAGON)
    return sp


def logo(s, x, y, w, dark=False):
    path = LOGO_DARK if dark else LOGO_COLOR
    ratio = 287.0 / 2400.0
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w),
                               height=Inches(w * ratio))
    return pic


def footer(s, dark=False, page=None):
    fg = TEXT_DIM_DARK if dark else TEXT_MUTED
    text(s, 0.55, 7.0, 8, 0.3, "Confidencial  ·  Perpetual Technologies © " + YEAR,
         size=8.5, color=fg, font="JetBrains Mono", anchor=MSO_ANCHOR.MIDDLE)
    if page is not None:
        text(s, 11.7, 7.0, 1.08, 0.3, str(page).zfill(2), size=8.5, color=fg,
             font="JetBrains Mono", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(s, label, title_runs, sub=None, page=None):
    """Header estandar de slide de contenido."""
    logo(s, 0.55, 0.5, 1.25)
    chip(s, 11.0, 0.52, "ESSENTIALS" if False else "PERPETUAL", fill=SURFACE2, fg=ACCENT, w=1.35)
    text(s, 0.55, 1.15, 9.5, 0.32, label, size=11, color=ACCENT, font=F_LABEL,
         spacing=1.4, upper=True)
    text(s, 0.5, 1.42, 11.2, 0.9, title_runs, size=30, font=F_DISPLAY, color=TEXT)
    if sub:
        text(s, 0.55, 2.28, 9.5, 0.5, sub, size=13, color=TEXT_MUTED, font=F_LIGHT)


def line_chart(s, x, y, w, h, cats, series, smooth=True, colors=None, dark=False):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    colors = colors or DATA_PALETTE
    for i, ser in enumerate(plot.series):
        ser.smooth = smooth
        lf = ser.format.line
        lf.color.rgb = C(colors[i % len(colors)])
        lf.width = Pt(3.0)
    _style_axes(chart, dark)
    return chart


def col_chart(s, x, y, w, h, cats, series, colors=None, dark=False, gap=80, overlap=-20):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = gap
    try:
        plot.overlap = overlap
    except Exception:
        pass
    colors = colors or DATA_PALETTE
    for i, ser in enumerate(plot.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = C(colors[i % len(colors)])
        ser.format.line.fill.background()
    _style_axes(chart, dark)
    return chart


def donut(s, x, y, size, pct, color=ACCENT, rest=SURFACE2, hole=60):
    cd = CategoryChartData()
    cd.categories = ["v", "r"]
    cd.add_series("d", (pct, 100 - pct))
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(size), Inches(size), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    try:
        plot.donut_hole_size = hole
    except Exception:
        pass
    pts = plot.series[0].points
    pts[0].format.fill.solid()
    pts[0].format.fill.fore_color.rgb = C(color)
    pts[1].format.fill.solid()
    pts[1].format.fill.fore_color.rgb = C(rest)
    for p in pts:
        p.format.line.color.rgb = C(BG)
        p.format.line.width = Pt(1.5)
    return chart


def _style_axes(chart, dark=False):
    axc = TEXT_DIM_DARK if dark else TEXT_MUTED
    try:
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(9)
        ca.tick_labels.font.name = F_BODY
        ca.tick_labels.font.color.rgb = C(axc)
        ca.format.line.color.rgb = C(BORDER_DARK if dark else BORDER)
        ca.major_tick_mark = XL_TICK_MARK.NONE
        ca.minor_tick_mark = XL_TICK_MARK.NONE
        ca.has_major_gridlines = False
    except Exception:
        pass
    try:
        va = chart.value_axis
        va.visible = False
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = C(BORDER_DARK if dark else BORDER)
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.color.rgb = C(axc)
        va.major_tick_mark = XL_TICK_MARK.NONE
    except Exception:
        pass


def card(s, x, y, w, h, fill=BG, radius=0.06, shadow=True, line=None):
    return rect(s, x, y, w, h, fill=fill, radius=radius, shadow=shadow, line=line)


def float_card(s, x, y, label, sub=None, w=2.85, icon=ACCENT, dark_icon=False,
               fg=TEXT, sub_fg=TEXT_MUTED, fill=BG):
    h = 0.92
    card(s, x, y, w, h, fill=fill, radius=0.16, shadow=True)
    hexagon(s, x + 0.18, y + 0.2, 0.52, fill=icon)
    rect(s, x + 0.30, y + 0.32, 0.28, 0.28, fill=BG if not dark_icon else SURFACE,
         shape=MSO_SHAPE.OVAL)
    tx = x + 0.92
    if sub:
        text(s, tx, y + 0.16, w - 1.0, 0.3, label, size=12, color=fg, font=F_HEAD)
        text(s, tx, y + 0.48, w - 1.0, 0.3, sub, size=10, color=sub_fg, font=F_BODY)
    else:
        text(s, tx, y, w - 1.0, h, label, size=12.5, color=fg, font=F_HEAD,
             anchor=MSO_ANCHOR.MIDDLE)


def stat(s, x, y, w, number, label, num_color=TEXT, lbl_color=TEXT_MUTED, num_size=34, align=PP_ALIGN.LEFT):
    text(s, x, y, w, 0.6, number, size=num_size, color=num_color, font=F_DISPLAY, align=align)
    text(s, x, y + num_size / 58.0, w, 0.35, label, size=11, color=lbl_color, font=F_LABEL,
         align=align, upper=True, spacing=0.6)


# ---------------------------------------------------------------------------
# Helpers especificos del template "Marketing"
# ---------------------------------------------------------------------------
def blob(s, x, y, d, fill):
    return rect(s, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)


def pill(s, x, y, w, label, fill=ACCENT, fg=TEXT_ON_DARK, arrow=True):
    rect(s, x, y, w, 0.62, fill=fill, radius=0.5, shadow=True)
    text(s, x + 0.34, y, w - 1.0, 0.62, label, size=11.5, color=fg, font=F_LABEL,
         anchor=MSO_ANCHOR.MIDDLE, upper=True, spacing=0.8)
    if arrow:
        rect(s, x + w - 0.74, y + 0.1, 0.42, 0.42, fill=BG, shape=MSO_SHAPE.OVAL)
        text(s, x + w - 0.74, y + 0.06, 0.42, 0.42, "›", size=17, color=fill, font=F_DISPLAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def photo_ph(s, x, y, w, h, radius=0.06, tint="E3ECFB"):
    """Placeholder de imagen: panel tintado + glifo de media."""
    rect(s, x, y, w, h, fill=tint, radius=radius)
    d = min(w, h) * 0.24
    cxp, cyp = x + w / 2, y + h / 2
    rect(s, cxp - d / 2, cyp - d / 2, d, d, fill=BG, shape=MSO_SHAPE.OVAL)
    rect(s, cxp - d * 0.16, cyp - d * 0.16, d * 0.32, d * 0.32, fill=ACCENT, shape=MSO_SHAPE.OVAL)


def tool_icon(s, x, y, d, fill):
    blob(s, x, y, d, fill)
    hexagon(s, x + d * 0.3, y + d * 0.3, d * 0.4, fill=BG)


def graphic(s, x, y, w, h, tint="E3ECFB", variant="abstract", shadow=False):
    """Grafico de marca en vez de foto (en Perpetual no usamos fotos de personas
    salvo en la slide de equipo)."""
    rect(s, x, y, w, h, fill=tint, radius=0.06, shadow=shadow)
    cx, cy = x + w / 2, y + h / 2
    if variant == "growth":
        n, bw, gap = 4, w * 0.13, w * 0.06
        total = n * bw + (n - 1) * gap
        bx, base = cx - total / 2, y + h * 0.8
        cols = [ACCENT, ACCENT2, BRAND_YELLOW, ACCENT]
        for i in range(n):
            bh = h * (0.16 + 0.13 * i)
            rect(s, bx + i * (bw + gap), base - bh, bw, bh, fill=cols[i], radius=0.2)
        rect(s, cx - w * 0.3, y + h * 0.16, h * 0.2, h * 0.2, fill=ACCENT, shape=MSO_SHAPE.OVAL)
        hexagon(s, cx + w * 0.16, y + h * 0.14, h * 0.16, fill=BRAND_YELLOW)
    elif variant == "quote":
        text(s, x, y + h * 0.02, w, h * 0.5, "“", size=80, color=ACCENT, font=F_DISPLAY, align=PP_ALIGN.CENTER)
        text(s, x, y + h * 0.62, w, h * 0.2, "★ ★ ★ ★ ★", size=16, color=BRAND_YELLOW,
             font=F_BODY, align=PP_ALIGN.CENTER)
    else:
        rect(s, cx - w * 0.28, cy - h * 0.16, h * 0.34, h * 0.34, fill=ACCENT, shape=MSO_SHAPE.OVAL)
        rect(s, cx + w * 0.03, cy - h * 0.02, h * 0.22, h * 0.22, fill=ACCENT2, shape=MSO_SHAPE.OVAL)
        rect(s, cx - w * 0.02, cy + h * 0.16, h * 0.13, h * 0.13, fill=BRAND_YELLOW, shape=MSO_SHAPE.OVAL)
        hexagon(s, cx + w * 0.12, cy - h * 0.26, h * 0.17, fill=BG)


def title(s, runs, x=0.7, y=0.7, w=7.5, size=33):
    logo(s, 0.6, 0.5, 1.15)
    text(s, x, y + 0.55, w, 1.2, runs, size=size, font=F_DISPLAY, color=TEXT, line_spacing=1.0)


# ===========================================================================
# 01 — Portada
# ===========================================================================
def m01():
    s = slide(BG)
    rect(s, 0, 4.55, 13.333, 2.95, fill=ACCENT, radius=0)
    blob(s, 9.0, -1.0, 3.2, BRAND_YELLOW)
    blob(s, 11.6, 2.2, 2.0, ACCENT2)
    blob(s, 8.0, 1.5, 3.3, ACCENT)
    hexagon(s, 8.55, 2.35, 2.1, fill=BG)
    rect(s, 9.25, 3.05, 0.7, 0.7, fill=ACCENT, shape=MSO_SHAPE.OVAL)
    blob(s, 11.4, 4.7, 0.8, BRAND_YELLOW)
    text(s, 12.3, 0.55, 0.7, 0.6, "+", size=30, color=ACCENT, font=F_DISPLAY)
    text(s, 7.4, 5.4, 0.6, 0.5, "+", size=24, color=ACCENT2, font=F_DISPLAY)
    logo(s, 0.7, 0.7, 1.5)
    text(s, 0.65, 1.5, 7.0, 2.0, [("Marketing para\n", TEXT, F_DISPLAY, 46), ("tu marca.", ACCENT, F_DISPLAY, 46)],
         line_spacing=0.98)
    text(s, 0.7, 3.5, 6.2, 0.3, "Estrategia y crecimiento", size=13, color=TEXT, font=F_HEAD, upper=False)
    text(s, 0.7, 3.85, 6.0, 0.8, "Convertimos tu marca en resultados medibles con data y creatividad.",
         size=12.5, color=TEXT_MUTED, font=F_BODY, line_spacing=1.3)
    pill(s, 0.7, 5.45, 3.0, "Empezar", fill=BG, fg=ACCENT)


# ===========================================================================
# 02 — Welcome
# ===========================================================================
def m02():
    s = slide(BG)
    title(s, [("Bienvenido a ", TEXT), ("Perpetual.", ACCENT)])
    blob(s, 5.2, 1.9, 1.7, ACCENT)
    graphic(s, 4.7, 1.9, 4.3, 3.0, tint="DBE7FB", variant="abstract")
    chip(s, 5.0, 4.55, "SOPORTE DE MARKETING", fill=ACCENT2, fg=TEXT_ON_DARK, w=2.7)
    stat(s, 0.7, 2.6, 3.5, "2.9M", "Impresiones gestionadas", num_color=ACCENT, num_size=40)
    text(s, 0.7, 3.9, 3.6, 2.5,
         "Acompanamos a tu equipo de principio a fin: estrategia, ejecucion y medicion en un solo lugar.",
         size=12.5, color=TEXT_MUTED, font=F_BODY, line_spacing=1.35)
    graphic(s, 9.4, 2.3, 3.3, 3.6, tint="FDE9D6", variant="abstract")
    footer(s, page=2)


# ===========================================================================
# 03 — Meet the team
# ===========================================================================
def m03():
    s = slide(BG)
    title(s, [("Conoce a nuestro\n", TEXT), ("equipo.", ACCENT)], w=5.0, size=30)
    people = [("Ana Rojas", "Marketing Lead"), ("Luis Vega", "Data Analyst"),
              ("Mara Diaz", "Creative Dir."), ("Iván Soto", "Growth")]
    for i, (n, role) in enumerate(people):
        x = 6.0 + i * 1.75
        photo_ph(s, x, 1.6, 1.55, 2.2, tint="DBE7FB" if i % 2 == 0 else "FDE9D6")
        text(s, x, 3.9, 1.65, 0.3, n, size=11.5, color=TEXT, font=F_HEAD)
        text(s, x, 4.2, 1.65, 0.3, role, size=9.5, color=ACCENT, font=F_LABEL, upper=True)
    blob(s, 11.9, 5.2, 1.0, BRAND_YELLOW)
    text(s, 0.7, 4.6, 4.8, 1.2,
         "Un equipo multidisciplinario enfocado en resultados, no en vanity metrics.",
         size=13, color=TEXT_MUTED, font=F_BODY, line_spacing=1.35)
    footer(s, page=3)


# ===========================================================================
# 04 — About current work
# ===========================================================================
def m04():
    s = slide(BG)
    title(s, [("Sobre el trabajo\n", TEXT), ("actual.", ACCENT)], size=30)
    blob(s, 6.7, 1.3, 3.0, ACCENT)
    graphic(s, 7.6, 1.5, 3.4, 2.4, tint="DBE7FB", variant="abstract")
    graphic(s, 6.2, 3.4, 3.4, 2.3, tint="FDE9D6", variant="abstract")
    blob(s, 9.7, 4.4, 1.9, ACCENT2)
    pill(s, 9.5, 4.9, 2.7, "Nueva via del exito", fill=ACCENT2, fg=TEXT_ON_DARK, arrow=False)
    text(s, 0.7, 2.8, 4.6, 1.8,
         "Cada cuenta tiene un plan vivo: lo medimos cada semana y ajustamos con base en datos.",
         size=13, color=TEXT_MUTED, font=F_BODY, line_spacing=1.35)
    text(s, 0.7, 4.7, 4.0, 0.3, "Auditoria y diagnostico", size=12, color=ACCENT, font=F_HEAD)
    text(s, 0.7, 5.1, 4.0, 0.3, "Plan de medios y contenido", size=12, color=ACCENT, font=F_HEAD)
    footer(s, page=4)


# ===========================================================================
# 05 — What we offer
# ===========================================================================
def m05():
    s = slide(BG)
    title(s, [("Lo que ", TEXT), ("ofrecemos.", ACCENT)])
    text(s, 8.2, 1.0, 4.4, 0.9, "Servicios end to end para que tu marca crezca con foco en retorno.",
         size=12.5, color=TEXT_MUTED, font=F_BODY, line_spacing=1.3)
    cards = [("SEO y contenido", ACCENT, TEXT_ON_DARK), ("Performance ads", BG_DARK, TEXT_ON_DARK),
             ("Brand y creatividad", ACCENT2, TEXT_ON_DARK)]
    for i, (t, fill, fg) in enumerate(cards):
        x = 0.7 + i * 4.1
        card(s, x, 2.6, 3.7, 3.9, fill=fill, radius=0.08, shadow=True)
        blob(s, x + 0.35, 2.95, 0.7, BG)
        hexagon(s, x + 0.5, 3.1, 0.4, fill=fill)
        text(s, x + 0.4, 4.0, 3.0, 0.7, t, size=15, color=fg, font=F_HEAD, line_spacing=1.05)
        text(s, x + 0.4, 4.9, 3.0, 1.4,
             "Estrategia, ejecucion y reporte claro de cada iniciativa.",
             size=11, color=fg, font=F_BODY, line_spacing=1.3)
    footer(s, page=5)


# ===========================================================================
# 06 — Market trends
# ===========================================================================
def m06():
    s = slide(BG)
    title(s, [("Tendencias de ", TEXT), ("mercado.", ACCENT)])
    blob(s, 3.0, 2.7, 2.9, ACCENT2)
    blob(s, 5.2, 3.2, 2.5, BRAND_YELLOW)
    text(s, 3.0, 3.7, 2.9, 0.7, "3.8M", size=30, color=TEXT_ON_DARK, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    text(s, 5.2, 4.0, 2.5, 0.6, "40.5K", size=26, color=TEXT, font=F_DISPLAY, align=PP_ALIGN.CENTER)
    text(s, 0.7, 3.4, 2.1, 0.8, "Alcance\nmensual", size=11, color=TEXT_MUTED, font=F_LABEL, upper=True, line_spacing=1.1)
    text(s, 8.1, 3.4, 2.4, 0.8, "Nuevos\nseguidores", size=11, color=TEXT_MUTED, font=F_LABEL, upper=True, line_spacing=1.1)
    col_chart(s, 9.4, 2.6, 3.3, 3.4, ["E", "F", "M", "A", "M"], [("a", [40, 55, 48, 70, 82])],
              colors=[BRAND_YELLOW], gap=55)
    footer(s, page=6)


# ===========================================================================
# 07 — Our solutions
# ===========================================================================
def m07():
    s = slide(BG)
    title(s, [("Nuestras ", TEXT), ("soluciones.", ACCENT)])
    text(s, 8.4, 1.0, 4.2, 0.9, "Dos frentes de trabajo que se refuerzan entre si.",
         size=12.5, color=TEXT_MUTED, font=F_BODY, line_spacing=1.3)
    rows = [("O1", "Adquisicion", "Campanas de performance optimizadas a costo por resultado."),
            ("O2", "Retencion", "Contenido y automatizaciones que sostienen el crecimiento.")]
    for i, (n, t, d) in enumerate(rows):
        y = 2.7 + i * 1.85
        card(s, 0.7, y, 11.9, 1.55, fill=ACCENT if i == 0 else BG_DARK, radius=0.08, shadow=True)
        text(s, 1.1, y + 0.4, 1.0, 0.8, n, size=26, color=TEXT_ON_DARK, font=F_DISPLAY)
        text(s, 2.3, y + 0.3, 3.0, 0.4, t, size=15, color=TEXT_ON_DARK, font=F_HEAD)
        text(s, 2.3, y + 0.75, 9.6, 0.6, d, size=12, color="DBE4FF", font=F_BODY, line_spacing=1.25)
    footer(s, page=7)


# ===========================================================================
# 08 — Marketing tools
# ===========================================================================
def m08():
    s = slide(BG)
    title(s, [("Herramientas de ", TEXT), ("marketing.", ACCENT)])
    tools = [("Audio social", ACCENT), ("SEO", ACCENT2), ("Video", BG_DARK),
             ("Social", BRAND_YELLOW), ("Automatizacion", ACCENT)]
    for i, (t, col) in enumerate(tools):
        x = 0.8 + i * 2.45
        tool_icon(s, x, 2.9, 1.4, col)
        text(s, x - 0.25, 4.5, 1.9, 0.6, t, size=11, color=TEXT, font=F_HEAD, align=PP_ALIGN.CENTER,
             line_spacing=1.05)
        text(s, x - 0.25, 5.15, 1.9, 0.5, "Stack y procesos", size=9.5, color=TEXT_MUTED,
             font=F_BODY, align=PP_ALIGN.CENTER)
    footer(s, page=8)


# ===========================================================================
# 09 — Customer reviews
# ===========================================================================
def m09():
    s = slide(BG)
    title(s, [("Opiniones de ", TEXT), ("clientes.", ACCENT)], w=4.5)
    text(s, 5.0, 1.05, 3.0, 0.8, "4,890 +", size=34, color=ACCENT, font=F_DISPLAY)
    graphic(s, 9.7, 1.0, 3.0, 3.4, tint="DBE7FB", variant="quote")
    blob(s, 12.2, 3.6, 1.0, BRAND_YELLOW)
    revs = [(ACCENT2, BG), (ACCENT, TEXT_ON_DARK)]
    for i, (icon, _) in enumerate(revs):
        y = 2.7 + i * 1.55
        fill = SURFACE if i == 0 else ACCENT
        fg = TEXT_DIM if i == 0 else TEXT_ON_DARK
        card(s, 0.7, y, 8.5, 1.35, fill=fill, radius=0.1, shadow=True, line=BORDER if i == 0 else None)
        blob(s, 1.0, y + 0.35, 0.65, icon)
        text(s, 1.95, y + 0.25, 7.0, 0.9,
             "El equipo entiende el negocio y entrega resultados claros mes a mes.",
             size=12, color=fg, font=F_BODY, line_spacing=1.25)
    footer(s, page=9)


# ===========================================================================
# 10 — New way of success
# ===========================================================================
def m10():
    s = slide(BG)
    title(s, [("Nueva via del ", TEXT), ("exito.", ACCENT)])
    graphic(s, 0.7, 2.7, 3.2, 3.6, tint="DBE7FB", variant="abstract")
    steps = [("O2", "Diagnostico", ACCENT, 4.4, 2.6),
             ("O3", "Estrategia", ACCENT2, 8.3, 2.6),
             ("O4", "Ejecucion", BRAND_YELLOW, 4.4, 4.5),
             ("O5", "Optimizacion", BG_DARK, 8.3, 4.5)]
    for n, t, col, x, y in steps:
        blob(s, x, y, 0.85, col)
        hexagon(s, x + 0.26, y + 0.26, 0.34, fill=BG)
        text(s, x + 1.05, y + 0.04, 1.0, 0.3, n, size=12, color=TEXT_MUTED, font=F_LABEL)
        text(s, x + 1.05, y + 0.32, 3.3, 0.35, t, size=14, color=TEXT, font=F_HEAD)
        text(s, x + 1.05, y + 0.72, 3.3, 0.4, "Paso del proceso con foco en dato.",
             size=10.5, color=TEXT_MUTED, font=F_BODY)
    footer(s, page=10)


# ===========================================================================
# 11 — Marketing strategy plan
# ===========================================================================
def m11():
    s = slide(BG)
    title(s, [("Plan de estrategia\n", TEXT), ("de marketing.", ACCENT)], size=30)
    blocks = [(ACCENT2, TEXT_ON_DARK, 7.0, 1.2, 5.2, "Objetivos y KPIs del trimestre"),
              (BG_DARK, TEXT_ON_DARK, 6.0, 2.5, 4.6, "Mensajes y audiencias clave"),
              (ACCENT, TEXT_ON_DARK, 7.4, 3.8, 5.0, "Mezcla de canales y presupuesto"),
              (SURFACE, TEXT_DIM, 6.4, 5.1, 5.4, "Calendario y responsables")]
    for fill, fg, x, y, w, t in blocks:
        card(s, x, y, w, 1.0, fill=fill, radius=0.12, shadow=True, line=BORDER if fill == SURFACE else None)
        text(s, x + 0.35, y, w - 0.7, 1.0, t, size=13, color=fg, font=F_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 2.9, 4.6, 2.0,
         "Un plan claro, por fases, con metas medibles y duenos definidos en cada etapa.",
         size=13, color=TEXT_MUTED, font=F_BODY, line_spacing=1.35)
    footer(s, page=11)


# ===========================================================================
# 12 — Grow with us
# ===========================================================================
def m12():
    s = slide(BG)
    title(s, [("Crece con ", TEXT), ("nosotros.", ACCENT)])
    card(s, 6.8, 1.7, 5.9, 4.6, fill=ACCENT, radius=0.1, shadow=True)
    text(s, 7.2, 2.1, 4.0, 0.8, "+58%", size=40, color=TEXT_ON_DARK, font=F_DISPLAY)
    text(s, 7.2, 3.0, 5.0, 0.4, "Crecimiento interanual promedio", size=12, color="DBE4FF", font=F_BODY)
    # linea de tendencia simple
    for i, (vx, vy) in enumerate([(7.4, 5.3), (8.6, 4.9), (9.8, 5.0), (11.0, 4.4), (12.0, 4.0)]):
        blob(s, vx, vy, 0.16, BRAND_YELLOW)
    stat(s, 0.7, 2.7, 3.5, "100K", "Leads generados", num_color=ACCENT, num_size=38)
    stat(s, 0.7, 4.3, 3.5, "0.2M", "Ventas atribuidas", num_color=TEXT, num_size=38)
    blob(s, 4.9, 1.5, 0.7, BRAND_YELLOW)
    footer(s, page=12)


# ===========================================================================
# 13 — Contact
# ===========================================================================
def m13():
    s = slide(BG)
    rect(s, 0, 0, 13.333, 3.0, fill=ACCENT, radius=0)
    blob(s, 10.8, -0.8, 2.2, BRAND_YELLOW)
    blob(s, 12.2, 1.8, 1.4, ACCENT2)
    logo(s, 0.7, 0.6, 1.5, dark=True)
    pill(s, 0.7, 1.6, 3.0, "Empezar", fill=BG, fg=ACCENT)
    text(s, 4.0, 1.55, 7.0, 0.8, "Hablemos de tu proxima campana.", size=15, color=TEXT_ON_DARK,
         font=F_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 3.7, 6.0, 1.0, [("Contac", TEXT, F_DISPLAY, 40), ("tanos.", ACCENT, F_DISPLAY, 40)])
    cols = [("Correo", "hola@perpetual.pe"), ("Telefono", "+51 999 999 999"), ("Web", "perpetual.pe")]
    for i, (t, v) in enumerate(cols):
        x = 0.7 + i * 4.1
        text(s, x, 5.4, 3.7, 0.3, t, size=11, color=ACCENT, font=F_LABEL, upper=True)
        text(s, x, 5.75, 3.7, 0.4, v, size=14, color=TEXT, font=F_HEAD)
    blob(s, 11.8, 5.6, 1.0, ACCENT)


for fn in [m01, m02, m03, m04, m05, m06, m07, m08, m09, m10, m11, m12, m13]:
    fn()

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("OK:", OUT, "| slides:", len(prs.slides._sldIdLst))

